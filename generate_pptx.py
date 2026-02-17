from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
TEAL = RGBColor(0x3D, 0xD5, 0xC0)
TEAL_BG = RGBColor(0xEA, 0xF8, 0xF6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
DARK_BG = RGBColor(0x3A, 0x3A, 0x3A)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

W = Inches(13.33)
H = Inches(7.5)


def add_rect(slide, l, t, w, h, fill_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_footer(slide, teal_color=True):
    c = TEAL if teal_color else RGBColor(0xFF, 0xFF, 0xFF)
    add_rect(slide, Inches(0.9), H - Inches(0.55), Inches(0.07), Inches(0.28), c)
    add_rect(slide, Inches(1.03), H - Inches(0.52), Inches(0.22), Inches(0.22), c)
    add_text(slide, "PASS", Inches(1.3), H - Inches(0.58), Inches(1.2), Inches(0.38),
             font_size=13, bold=True, color=c)


def slide_header(slide, label, title, subtitle=None):
    add_rect(slide, Inches(0.9), Inches(0.5), Inches(0.07), Inches(0.34), TEAL)
    add_rect(slide, Inches(1.03), Inches(0.54), Inches(0.34), Inches(0.34), TEAL)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.38), Inches(11.3), Inches(1.6))
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = label
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = TEAL

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = GRAY


# ==================== SLIDE 01: TITLE ====================
slide1 = prs.slides.add_slide(blank_layout)
add_rect(slide1, 0, 0, W, H, TEAL)

add_text(slide1, "デザイン定額サービス", Inches(0.9), Inches(1.4), Inches(6), Inches(0.5),
         font_size=16, bold=True, color=WHITE)
add_rect(slide1, Inches(0.9), Inches(2.0), Inches(0.13), Inches(0.62), WHITE)
add_rect(slide1, Inches(1.1), Inches(2.0), Inches(0.62), Inches(0.62), WHITE)
add_text(slide1, "PASS", Inches(1.85), Inches(1.85), Inches(5), Inches(0.9),
         font_size=62, bold=True, color=WHITE)

add_text(slide1, "デザインサブスクリプション「PASS」", Inches(0.9), Inches(3.2), Inches(11), Inches(0.85),
         font_size=38, bold=True, color=WHITE)
add_text(slide1, "採用コスト0で、即戦力デザインチームをあなたの手元に。", Inches(0.9), Inches(4.2), Inches(10), Inches(0.6),
         font_size=22, color=WHITE)
add_text(slide1, "株式会社SAKAZUKI", Inches(0.9), Inches(5.1), Inches(6), Inches(0.5),
         font_size=16, color=WHITE)

# Deco
add_rect(slide1, W - Inches(3.2), H - Inches(3.2), Inches(3.5), Inches(3.5), RGBColor(0xFF, 0xFF, 0xFF))

# ==================== SLIDE 02: COMPANY ====================
slide2 = prs.slides.add_slide(blank_layout)
add_rect(slide2, 0, 0, W, H, WHITE)
slide_header(slide2, "About Us", "私たちは「志」をデザインで具現化するパートナーです")
slide_footer(slide2)

col_w = Inches(3.6)
# Col1: Vision
add_rect(slide2, Inches(0.9), Inches(2.0), col_w, Inches(0.04), TEAL)
add_text(slide2, "Vision", Inches(0.9), Inches(2.1), col_w, Inches(0.4), font_size=15, bold=True, color=TEAL)
add_text(slide2, "SAKAZUKIのビジョンを一言で", Inches(0.9), Inches(2.65), col_w, Inches(0.5), font_size=16, color=DARK)

# Col2: Company
add_rect(slide2, Inches(4.7), Inches(2.0), col_w, Inches(0.04), TEAL)
txb2 = slide2.shapes.add_textbox(Inches(4.7), Inches(2.1), col_w, Inches(3.2))
tf2 = txb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "Company"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = TEAL
add_para(tf2, "株式会社SAKAZUKI", 16, color=DARK)
add_para(tf2, "代表取締役 [代表者名]", 14, color=GRAY)
add_para(tf2, "", 8)
add_para(tf2, "Business", 15, bold=True, color=TEAL)
add_para(tf2, "・デザインサブスクリプション事業「PASS」", 14, color=DARK)
add_para(tf2, "・ブランディング・デザイン制作", 14, color=DARK)

# Col3: 実績
add_rect(slide2, Inches(8.5), Inches(2.0), col_w, Inches(0.04), TEAL)
txb3 = slide2.shapes.add_textbox(Inches(8.5), Inches(2.1), col_w, Inches(3.2))
tf3 = txb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
r3 = p3.add_run(); r3.text = "実績"; r3.font.size = Pt(15); r3.font.bold = True; r3.font.color.rgb = TEAL
for item in ["累計導入企業 [XX]社以上", "継続率 [XX]%", "平均契約期間 [XX]ヶ月", "業種を問わず幅広い実績"]:
    add_para(tf3, "・" + item, 14, color=DARK)


# ==================== SLIDE 03: PROBLEM ====================
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, W, H, WHITE)
slide_header(slide3, "Problem", "課題提起：クリエイティブの「3つの壁」")
slide_footer(slide3)

problems = [
    ("WALL 01", "採用難易度の高騰", "有効求人倍率 [X.X]倍\n採用期間 [3〜6]ヶ月", "採れない・育てられない"),
    ("WALL 02", "外注管理の複雑さ", "見積もり、契約、ディレクションで\nコア業務を圧迫", "手間・時間・コストが爆増"),
    ("WALL 03", "品質のバラつき", "クラウドソーシング等の品質不安定\nブランド毀損リスク", "安心して任せられない"),
]
for i, (num, h3, body, accent) in enumerate(problems):
    x = Inches(0.9 + i * 4.15)
    add_rect(slide3, x, Inches(1.9), Inches(3.9), Inches(5.0), TEAL_BG)
    add_rect(slide3, x + Inches(0.25), Inches(2.1), Inches(1.6), Inches(0.33), TEAL)
    add_text(slide3, num, x + Inches(0.25), Inches(2.08), Inches(1.6), Inches(0.36),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide3, h3, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.65),
             font_size=22, bold=True, color=DARK)
    add_text(slide3, body, x + Inches(0.3), Inches(3.4), Inches(3.3), Inches(0.9),
             font_size=15, color=GRAY)
    add_text(slide3, accent, x + Inches(0.3), Inches(4.5), Inches(3.3), Inches(0.5),
             font_size=16, bold=True, color=ORANGE)


# ==================== SLIDE 04: SOLUTION ====================
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, W, H, WHITE)
slide_header(slide4, "Solution", "解決策：デザインサブスク「PASS」とは", "必要な時に、必要なだけ。月額制デザインチーム。")
slide_footer(slide4)

add_text(slide4, "採用・育成・管理コストをゼロにし、ビジネスを加速させます。",
         Inches(0.9), Inches(2.1), Inches(6.2), Inches(0.75), font_size=21, bold=True, color=DARK)
add_text(slide4, "PASSは月額固定でデザインチームがそのままあなたの会社の\n制作部門として機能します。依頼はチャット一本で完結。",
         Inches(0.9), Inches(3.0), Inches(6.2), Inches(1.0), font_size=16, color=GRAY)

# Before
add_rect(slide4, Inches(7.3), Inches(2.0), Inches(5.6), Inches(1.75), LIGHT_GRAY_BG)
add_rect(slide4, Inches(7.3), Inches(2.0), Inches(0.07), Inches(1.75), RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide4, "■ Before", Inches(7.5), Inches(2.12), Inches(5.2), Inches(0.4), font_size=13, bold=True, color=GRAY)
add_text(slide4, "✕ 採用できない\n✕ 見積もりが面倒\n✕ クオリティ不安定",
         Inches(7.5), Inches(2.55), Inches(5.2), Inches(1.0), font_size=14, color=GRAY)

# After
add_rect(slide4, Inches(7.3), Inches(3.9), Inches(5.6), Inches(1.75), TEAL_BG)
add_rect(slide4, Inches(7.3), Inches(3.9), Inches(0.07), Inches(1.75), TEAL)
add_text(slide4, "■ After (PASS)", Inches(7.5), Inches(4.02), Inches(5.2), Inches(0.4), font_size=13, bold=True, color=TEAL)
add_text(slide4, "✓ チャット一本でプロに依頼\n✓ 月額固定・追加費用なし\n✓ 厳選された専任チーム",
         Inches(7.5), Inches(4.45), Inches(5.2), Inches(1.0), font_size=14, bold=True, color=DARK)


# ==================== SLIDE 05: STRENGTH 1 ====================
slide5 = prs.slides.add_slide(blank_layout)
add_rect(slide5, 0, 0, W, H, WHITE)
slide_header(slide5, "Strength 01", "PASSの強み①：広範囲な対応領域", "Webも、紙も、映像も。あらゆる制作物をワンストップで。")
slide_footer(slide5)

add_text(slide5, "PASSなら、あらゆるデザインニーズに対応。複数の外注先を管理する手間がゼロに。",
         Inches(0.9), Inches(2.1), Inches(6.2), Inches(0.75), font_size=17, bold=True, color=DARK)

for i, (bullet, text) in enumerate([
    ("W", "Web：LPデザイン / バナー / UIデザイン"),
    ("G", "Graphic：チラシ / パンフレット / 名刺 / ロゴ"),
    ("B", "Business：営業資料 / ホワイトペーパー"),
    ("M", "Movie：動画編集（※プランによる）"),
]):
    y = Inches(3.0 + i * 0.7)
    add_rect(slide5, Inches(0.9), y, Inches(0.3), Inches(0.3), TEAL)
    add_text(slide5, bullet, Inches(0.9), y - Inches(0.02), Inches(0.3), Inches(0.34), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide5, text, Inches(1.3), y, Inches(5.4), Inches(0.4), font_size=16, color=DARK)

for i, (title, body) in enumerate([
    ("対応制作物 20種類以上", "名刺・チラシ・ポスター・パンフレット・メニュー・バナー・ロゴ・キービジュアル・Webデザイン・Webコーディング・動画など多数。"),
    ("月額依頼し放題", "定額制だからコストを気にせず依頼できます。チケット制で優先度も柔軟に調整可能。"),
]):
    y = Inches(2.0 + i * 2.55)
    add_rect(slide5, Inches(7.3), y, Inches(5.6), Inches(2.3), TEAL_BG)
    add_text(slide5, title, Inches(7.6), y + Inches(0.3), Inches(5.1), Inches(0.5), font_size=14, bold=True, color=TEAL)
    add_text(slide5, body, Inches(7.6), y + Inches(0.85), Inches(5.1), Inches(1.2), font_size=14, color=DARK)


# ==================== SLIDE 06: STRENGTH 2 ====================
slide6 = prs.slides.add_slide(blank_layout)
add_rect(slide6, 0, 0, W, H, WHITE)
slide_header(slide6, "Strength 02", "PASSの強み②：圧倒的なスピードと質", "ビジネスを止めないスピード。最短1営業日で初稿提出。")
slide_footer(slide6)

add_text(slide6, "専任ディレクターによるダブルチェック体制で、修正の手間を最小限に。",
         Inches(0.9), Inches(2.1), Inches(6.2), Inches(0.6), font_size=17, bold=True, color=DARK)

for i, (bullet, text) in enumerate([
    ("⚡", "バナー：[1〜2]営業日"),
    ("⚡", "チラシ：[2〜3]営業日"),
    ("⚡", "LP：[5〜7]営業日"),
]):
    y = Inches(3.0 + i * 0.72)
    add_rect(slide6, Inches(0.9), y, Inches(0.3), Inches(0.3), TEAL)
    add_text(slide6, bullet, Inches(0.9), y - Inches(0.02), Inches(0.3), Inches(0.34), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide6, text, Inches(1.3), y, Inches(5.4), Inches(0.4), font_size=16, color=DARK)

for i, (title, body) in enumerate([
    ("専任ディレクター制", "ディレクターがあなたの会社を深く理解した上で制作をアサイン。品質とスピードを両立します。"),
    ("修正無制限", "追加料金なしで何度でも修正対応。「思っていたのと違う…」がなくなります。"),
]):
    y = Inches(2.0 + i * 2.55)
    add_rect(slide6, Inches(7.3), y, Inches(5.6), Inches(2.3), TEAL_BG)
    add_text(slide6, title, Inches(7.6), y + Inches(0.3), Inches(5.1), Inches(0.5), font_size=14, bold=True, color=TEAL)
    add_text(slide6, body, Inches(7.6), y + Inches(0.85), Inches(5.1), Inches(1.2), font_size=14, color=DARK)


# ==================== SLIDE 07: STRENGTH 3 ====================
slide7 = prs.slides.add_slide(blank_layout)
add_rect(slide7, 0, 0, W, H, WHITE)
slide_header(slide7, "Strength 03", "PASSの強み③：事業理解に基づくデザイン", "SAKAZUKIの強みである「構造思考」を活用。")
slide_footer(slide7)

add_text(slide7, "単なる作業代行ではなく、貴社の事業課題やKPIを理解した上で、\n「売れる」「伝わる」デザインを設計します。",
         Inches(0.9), Inches(2.1), Inches(6.2), Inches(1.0), font_size=17, bold=True, color=DARK)

for i, (bullet, text) in enumerate([
    ("★", "経営者視点での壁打ちも可能"),
    ("★", "ブランドガイドラインの構築・維持"),
    ("★", "KPIドリブンなクリエイティブ改善"),
]):
    y = Inches(3.3 + i * 0.72)
    add_rect(slide7, Inches(0.9), y, Inches(0.3), Inches(0.3), TEAL)
    add_text(slide7, bullet, Inches(0.9), y - Inches(0.02), Inches(0.3), Inches(0.34), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide7, text, Inches(1.3), y, Inches(5.4), Inches(0.4), font_size=16, color=DARK)

for i, (title, body) in enumerate([
    ("構造思考デザイン", "伝えたいことを整理し、ターゲットに最も刺さる構成と表現を提案。デザインを「コスト」から「投資」に変えます。"),
    ("継続的なブランド強化", "月次での振り返りとアップデートで、御社のブランド価値を継続的に高めていきます。"),
]):
    y = Inches(2.0 + i * 2.55)
    add_rect(slide7, Inches(7.3), y, Inches(5.6), Inches(2.3), TEAL_BG)
    add_text(slide7, title, Inches(7.6), y + Inches(0.3), Inches(5.1), Inches(0.5), font_size=14, bold=True, color=TEAL)
    add_text(slide7, body, Inches(7.6), y + Inches(0.85), Inches(5.1), Inches(1.2), font_size=14, color=DARK)


# ==================== SLIDE 08: FLOW ====================
slide8 = prs.slides.add_slide(blank_layout)
add_rect(slide8, 0, 0, W, H, WHITE)
slide_header(slide8, "How it Works", "制作フロー：チャット一本で完結", "まるで「隣の席」にいるような感覚で。")
slide_footer(slide8)

steps = [
    ("STEP 01", "Request", "チャットでテンプレに沿って依頼内容を記入。打ち合わせや電話は原則不要。"),
    ("STEP 02", "Direction", "専任ディレクターが内容を確認し、最適なデザイナーをアサイン。"),
    ("STEP 03", "Design", "制作開始。初稿を提出、お客様に確認いただきます。"),
    ("STEP 04", "Check", "初稿提出・修正対応。追加料金なしで何度でも修正可能。"),
    ("STEP 05", "納品", "最終確認後、納品データをお渡し。次の依頼もすぐに開始できます。"),
]
step_w = Inches(2.35)
for i, (num, name, body) in enumerate(steps):
    x = Inches(0.75 + i * 2.5)
    add_rect(slide8, x, Inches(2.0), step_w, Inches(0.85), TEAL)
    add_text(slide8, num, x, Inches(2.03), step_w, Inches(0.28), font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide8, name, x, Inches(2.3), step_w, Inches(0.45), font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide8, x + Inches(0.08), Inches(2.85), step_w - Inches(0.16), Inches(3.4), TEAL_BG)
    add_text(slide8, body, x + Inches(0.2), Inches(2.95), step_w - Inches(0.3), Inches(3.15), font_size=13, color=DARK)

add_text(slide8, "※面倒なMTGや電話は原則不要", Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
         font_size=15, color=GRAY, align=PP_ALIGN.CENTER)


# ==================== SLIDE 09: MENU ====================
slide9 = prs.slides.add_slide(blank_layout)
add_rect(slide9, 0, 0, W, H, WHITE)
slide_header(slide9, "Service Menu", "提供メニュー・チケット制")
slide_footer(slide9)

add_text(slide9, "1チケット（工数）で、これだけの制作が可能です。\nプランにより毎月の付与数が異なります。",
         Inches(0.9), Inches(1.9), Inches(5.6), Inches(0.9), font_size=16, color=DARK)

add_rect(slide9, Inches(0.9), Inches(2.95), Inches(5.6), Inches(1.8), TEAL)
add_text(slide9, "チケット制", Inches(0.9), Inches(3.25), Inches(5.6), Inches(0.7),
         font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide9, "プランに応じた月間チケットで柔軟にデザインを発注",
         Inches(0.9), Inches(3.98), Inches(5.6), Inches(0.5), font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

menu_items = [
    ("バナー制作", "0.5 チケット"),
    ("名刺デザイン", "1.0 チケット"),
    ("A4チラシ（片面）", "1.5 チケット"),
    ("A4チラシ（両面）", "2.0 チケット"),
    ("キービジュアル", "2.5 チケット"),
    ("ロゴ作成", "3.5 チケット"),
    ("LPデザイン（SP）", "5.0 チケット"),
    ("営業資料作成", "2.0 チケット"),
]
for i, (name, ticket) in enumerate(menu_items):
    row = i // 2
    col = i % 2
    x = Inches(7.1 + col * 3.1)
    y = Inches(1.9 + row * 1.1)
    add_rect(slide9, x, y, Inches(2.9), Inches(0.9), TEAL_BG)
    add_text(slide9, name, x + Inches(0.15), y + Inches(0.15), Inches(1.7), Inches(0.5), font_size=14, bold=True, color=DARK)
    add_text(slide9, ticket, x + Inches(1.85), y + Inches(0.15), Inches(0.95), Inches(0.5), font_size=13, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)


# ==================== SLIDE 10: CASE STUDY A ====================
slide10 = prs.slides.add_slide(blank_layout)
add_rect(slide10, 0, 0, W, H, WHITE)
slide_header(slide10, "Case Study 01", "導入事例①：株式会社[企業名A] 様")
slide_footer(slide10)

add_text(slide10, "デザイン外注から切り替え。コスト [XX]% 削減に成功。",
         Inches(0.9), Inches(1.95), Inches(6.6), Inches(0.6), font_size=19, bold=True, color=DARK)

add_rect(slide10, Inches(0.9), Inches(2.65), Inches(6.6), Inches(1.35), LIGHT_GRAY_BG)
add_rect(slide10, Inches(0.9), Inches(2.65), Inches(0.07), Inches(1.35), RGBColor(0xDD, 0xDD, 0xDD))
add_text(slide10, "■ Before", Inches(1.1), Inches(2.75), Inches(6.2), Inches(0.35), font_size=13, bold=True, color=GRAY)
add_text(slide10, "都度外注でコストが積み重なり、担当者の管理工数も膨大に。品質もバラバラで社内からの不満も。",
         Inches(1.1), Inches(3.15), Inches(6.2), Inches(0.7), font_size=13, color=GRAY)

add_rect(slide10, Inches(0.9), Inches(4.1), Inches(6.6), Inches(1.35), TEAL_BG)
add_rect(slide10, Inches(0.9), Inches(4.1), Inches(0.07), Inches(1.35), TEAL)
add_text(slide10, "■ After", Inches(1.1), Inches(4.2), Inches(6.2), Inches(0.35), font_size=13, bold=True, color=TEAL)
add_text(slide10, "PASSに切り替え後、月々のデザインコストが大幅削減。担当者の工数もほぼゼロに。",
         Inches(1.1), Inches(4.6), Inches(6.2), Inches(0.7), font_size=13, color=DARK)

add_rect(slide10, Inches(0.9), Inches(5.55), Inches(6.6), Inches(1.1), TEAL)
add_text(slide10, "「月次コストが[XX]万円→[XX]万円に。さらにデザインのスピードも上がり、\n本業に集中できるようになりました」",
         Inches(1.1), Inches(5.65), Inches(6.2), Inches(0.9), font_size=13, color=WHITE)

add_rect(slide10, Inches(8.1), Inches(1.95), Inches(4.7), Inches(2.2), TEAL)
add_text(slide10, "53%", Inches(8.1), Inches(2.3), Inches(4.7), Inches(1.2), font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide10, "コスト削減達成", Inches(8.1), Inches(3.3), Inches(4.7), Inches(0.5), font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide10, Inches(8.1), Inches(4.35), Inches(4.7), Inches(2.2), TEAL_BG)
add_text(slide10, "11万円", Inches(8.1), Inches(4.7), Inches(4.7), Inches(1.0), font_size=50, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide10, "月間コスト削減額", Inches(8.1), Inches(5.7), Inches(4.7), Inches(0.5), font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ==================== SLIDE 11: CASE STUDY B ====================
slide11 = prs.slides.add_slide(blank_layout)
add_rect(slide11, 0, 0, W, H, WHITE)
slide_header(slide11, "Case Study 02", "導入事例②：株式会社[企業名B] 様")
slide_footer(slide11)

add_text(slide11, "採用資料の刷新で、エントリー数が昨対比 [2] 倍に急増。",
         Inches(0.9), Inches(1.95), Inches(6.6), Inches(0.6), font_size=19, bold=True, color=DARK)

add_rect(slide11, Inches(0.9), Inches(2.65), Inches(6.6), Inches(1.35), LIGHT_GRAY_BG)
add_rect(slide11, Inches(0.9), Inches(2.65), Inches(0.07), Inches(1.35), RGBColor(0xDD, 0xDD, 0xDD))
add_text(slide11, "■ Before", Inches(1.1), Inches(2.75), Inches(6.2), Inches(0.35), font_size=13, bold=True, color=GRAY)
add_text(slide11, "資料デザインが古い、SNSトンマナばらつき。優秀層からのエントリーが減少傾向に。",
         Inches(1.1), Inches(3.15), Inches(6.2), Inches(0.7), font_size=13, color=GRAY)

add_rect(slide11, Inches(0.9), Inches(4.1), Inches(6.6), Inches(1.35), TEAL_BG)
add_rect(slide11, Inches(0.9), Inches(4.1), Inches(0.07), Inches(1.35), TEAL)
add_text(slide11, "■ After", Inches(1.1), Inches(4.2), Inches(6.2), Inches(0.35), font_size=13, bold=True, color=TEAL)
add_text(slide11, "トンマナ統一、資料フルリニューアル、優秀層からの応募増。採用広報の工数も大幅削減。",
         Inches(1.1), Inches(4.6), Inches(6.2), Inches(0.7), font_size=13, color=DARK)

add_rect(slide11, Inches(0.9), Inches(5.55), Inches(6.6), Inches(1.1), TEAL)
add_text(slide11, "「採用広報の工数が減り、応募の質が上がりました。デザインへの投資対効果が非常に高い」",
         Inches(1.1), Inches(5.65), Inches(6.2), Inches(0.9), font_size=13, color=WHITE)

add_rect(slide11, Inches(8.1), Inches(1.95), Inches(4.7), Inches(2.2), TEAL)
add_text(slide11, "2倍", Inches(8.1), Inches(2.3), Inches(4.7), Inches(1.2), font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide11, "エントリー数 昨対比", Inches(8.1), Inches(3.3), Inches(4.7), Inches(0.5), font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide11, Inches(8.1), Inches(4.35), Inches(4.7), Inches(2.2), TEAL_BG)
add_text(slide11, "90%", Inches(8.1), Inches(4.7), Inches(4.7), Inches(1.0), font_size=60, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide11, "ディレクション工数削減", Inches(8.1), Inches(5.7), Inches(4.7), Inches(0.5), font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ==================== SLIDE 12: PORTFOLIO ====================
slide12 = prs.slides.add_slide(blank_layout)
add_rect(slide12, 0, 0, W, H, WHITE)
slide_header(slide12, "Portfolio", "Design Portfolio", "多様なテイストに対応可能。")
slide_footer(slide12)

portfolio_items = [
    ("🌐", "Webサイト"), ("🎨", "バナー"), ("📄", "パンフレット"), ("✨", "イラスト"),
    ("📱", "SNS投稿"), ("🏷️", "ロゴ"), ("📊", "営業資料"), ("📽️", "動画編集"),
]
for i, (icon, label) in enumerate(portfolio_items):
    row = i // 4
    col = i % 4
    x = Inches(0.8 + col * 3.15)
    y = Inches(1.95 + row * 2.25)
    add_rect(slide12, x, y, Inches(3.0), Inches(2.0), TEAL_BG)
    add_text(slide12, icon, x, y + Inches(0.3), Inches(3.0), Inches(0.75), font_size=34, align=PP_ALIGN.CENTER)
    add_text(slide12, label, x, y + Inches(1.15), Inches(3.0), Inches(0.5), font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide12, "※実績画像は別途ご覧いただけます", Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
         font_size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ==================== SLIDE 13: QUALITY ====================
slide13 = prs.slides.add_slide(blank_layout)
add_rect(slide13, 0, 0, W, H, WHITE)
slide_header(slide13, "Quality & Security", "品質・体制への安心")
slide_footer(slide13)

for i, (icon, title, body) in enumerate([
    ("🎯", "厳選採用", "採用合格率 [X]%。実務経験豊富なプロのデザイナーのみ在籍。コンペや審査を通過したメンバーで構成。"),
    ("👥", "チーム体制", "ディレクターを含めたチームで品質管理。担当デザイナー + ディレクターによるダブルチェック体制。"),
    ("🔒", "セキュリティ", "全スタッフとNDA締結済み。機密情報も安心してお取り扱いいただけます。"),
]):
    x = Inches(0.9 + i * 4.15)
    add_rect(slide13, x, Inches(1.95), Inches(3.9), Inches(4.6), TEAL_BG)
    add_text(slide13, icon, x, Inches(2.25), Inches(3.9), Inches(0.85), font_size=42, align=PP_ALIGN.CENTER)
    add_text(slide13, title, x, Inches(3.2), Inches(3.9), Inches(0.6), font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide13, body, x + Inches(0.3), Inches(3.95), Inches(3.3), Inches(2.2), font_size=14, color=DARK)


# ==================== SLIDE 14: COST SIM 1 ====================
slide14 = prs.slides.add_slide(blank_layout)
add_rect(slide14, 0, 0, W, H, WHITE)
slide_header(slide14, "Cost Simulation 01", "コストシミュレーション①（対 採用）")
slide_footer(slide14)

add_rect(slide14, Inches(0.9), Inches(1.9), Inches(5.6), Inches(2.35), DARK_BG)
txb14a = slide14.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(2.1))
tf14a = txb14a.text_frame; tf14a.word_wrap = True
p = tf14a.paragraphs[0]; r = p.add_run(); r.text = "■A社（正社員採用）"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
add_para(tf14a, "年収 [500]万 / 採用費 [150]万 / 設備・福利厚生 等", 13, color=WHITE)
add_para(tf14a, "¥750万円 / 年間コスト（概算）", 20, bold=True, color=WHITE)

add_rect(slide14, Inches(0.9), Inches(4.35), Inches(5.6), Inches(2.35), TEAL)
txb14b = slide14.shapes.add_textbox(Inches(1.1), Inches(4.45), Inches(5.1), Inches(2.1))
tf14b = txb14b.text_frame; tf14b.word_wrap = True
p = tf14b.paragraphs[0]; r = p.add_run(); r.text = "■PASS（スタンダード）"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
add_para(tf14b, "月額 [20]万円 × 12ヶ月 / 採用リスクなし / 解約リスクなし", 13, color=WHITE)
add_para(tf14b, "¥240万円 / 年間コスト", 20, bold=True, color=WHITE)

add_text(slide14, "→", Inches(6.8), Inches(3.3), Inches(0.8), Inches(0.8), font_size=48, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(slide14, Inches(7.4), Inches(2.8), Inches(1.9), Inches(1.9), ORANGE)
add_text(slide14, "[1/3]", Inches(7.4), Inches(3.0), Inches(1.9), Inches(0.9), font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide14, "以下に削減", Inches(7.4), Inches(3.88), Inches(1.9), Inches(0.45), font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide14, "コスト削減 + 解約リスクなし", Inches(6.6), Inches(5.1), Inches(4.0), Inches(0.5), font_size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide14, "採用コスト・教育コスト・設備コストを含めると\n正社員雇用は年間750万円以上。\nPASSなら1/3以下でプロチームが動きます。",
         Inches(9.7), Inches(2.8), Inches(3.3), Inches(2.5), font_size=14, color=DARK)


# ==================== SLIDE 15: COST SIM 2 ====================
slide15 = prs.slides.add_slide(blank_layout)
add_rect(slide15, 0, 0, W, H, WHITE)
slide_header(slide15, "Cost Simulation 02", "コストシミュレーション②（対 外注）", "「見えないコスト」であるディレクション工数を [90]% 削減。")
slide_footer(slide15)

add_rect(slide15, Inches(0.9), Inches(2.1), Inches(5.6), Inches(2.2), DARK_BG)
txb15a = slide15.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5.1), Inches(2.0))
tf15a = txb15a.text_frame; tf15a.word_wrap = True
p = tf15a.paragraphs[0]; r = p.add_run(); r.text = "■外注（従来）"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
add_para(tf15a, "業者選定・見積もり・契約 / MTG・ディレクション", 13, color=WHITE)
add_para(tf15a, "月[20]時間 / 管理工数", 20, bold=True, color=WHITE)

add_rect(slide15, Inches(0.9), Inches(4.45), Inches(5.6), Inches(2.1), TEAL)
txb15b = slide15.shapes.add_textbox(Inches(1.1), Inches(4.55), Inches(5.1), Inches(1.9))
tf15b = txb15b.text_frame; tf15b.word_wrap = True
p = tf15b.paragraphs[0]; r = p.add_run(); r.text = "■PASS"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
add_para(tf15b, "チャットで依頼するだけ / = 月間 [2]時間", 13, color=WHITE)
add_para(tf15b, "月[2]時間 / 管理工数", 20, bold=True, color=WHITE)

add_text(slide15, "→", Inches(6.8), Inches(3.5), Inches(0.8), Inches(0.8), font_size=48, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(slide15, Inches(7.4), Inches(3.0), Inches(1.9), Inches(1.9), ORANGE)
add_text(slide15, "90%", Inches(7.4), Inches(3.2), Inches(1.9), Inches(0.9), font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide15, "工数削減", Inches(7.4), Inches(4.0), Inches(1.9), Inches(0.45), font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide15, "空いた時間をコア業務へ集中", Inches(6.6), Inches(5.2), Inches(4.0), Inches(0.5), font_size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide15, "管理工数の削減により、\n担当者はコア業務に集中できます。\n月18時間の削減 = 年間216時間の創出。",
         Inches(9.7), Inches(3.0), Inches(3.3), Inches(2.2), font_size=14, color=DARK)


# ==================== SLIDE 16: PRICING ====================
slide16 = prs.slides.add_slide(blank_layout)
add_rect(slide16, 0, 0, W, H, WHITE)
slide_header(slide16, "Plan & Pricing", "料金プラン")
slide_footer(slide16)

plans = [
    ("Light", "ライト", "¥[X]万円/月", "バナー制作メイン", LIGHT_GRAY_BG, DARK,
     TEAL, ["月[X]チケット付与", "Web/Graphicデザイン対応", "修正無制限", "チャットサポート"], False),
    ("Standard", "スタンダード", "¥[Y]万円/月", "LP・資料作成など", TEAL, WHITE,
     WHITE, ["月[Y]チケット付与", "全ジャンルのデザイン対応", "修正無制限", "専任ディレクター制", "Webコーディング対応"], True),
    ("Premium", "プレミアム", "¥[Z]万円/月", "動画対応・最優先対応", NEAR_BLACK, WHITE,
     RGBColor(0xFF, 0xD7, 0x00), ["月[Z]チケット付与", "動画編集対応", "最優先対応", "月次戦略MTG", "ブランドガイドライン策定"], False),
]
for i, (badge, name, price, note, bg, text_color, check_color, features, recommended) in enumerate(plans):
    x = Inches(0.7 + i * 4.22)
    add_rect(slide16, x, Inches(1.75), Inches(4.0), Inches(5.55), bg)
    if recommended:
        add_text(slide16, "★推奨", x + Inches(2.6), Inches(1.92), Inches(1.25), Inches(0.38),
                 font_size=12, bold=True, color=TEAL)
    add_text(slide16, badge, x + Inches(0.3), Inches(1.95), Inches(2.0), Inches(0.4), font_size=12, bold=True, color=check_color)
    add_text(slide16, name, x + Inches(0.3), Inches(2.45), Inches(3.4), Inches(0.6), font_size=26, bold=True, color=text_color)
    add_text(slide16, price, x + Inches(0.3), Inches(3.08), Inches(3.4), Inches(0.65), font_size=28, bold=True, color=text_color)
    add_text(slide16, note, x + Inches(0.3), Inches(3.75), Inches(3.4), Inches(0.4), font_size=13, color=GRAY if bg == LIGHT_GRAY_BG else text_color)
    for j, feat in enumerate(features):
        y = Inches(4.25 + j * 0.48)
        add_text(slide16, "✓ " + feat, x + Inches(0.3), y, Inches(3.4), Inches(0.42), font_size=13, color=text_color)


# ==================== SLIDE 17: CONTRACT ====================
slide17 = prs.slides.add_slide(blank_layout)
add_rect(slide17, 0, 0, W, H, WHITE)
slide_header(slide17, "Getting Started", "契約の流れ・トライアル", "まずは [1ヶ月] から。リスクなく始められます。")
slide_footer(slide17)

for i, (num, h4, p_text) in enumerate([
    ("1", "お問い合わせ", "フォームまたはメールよりご連絡ください"),
    ("2", "無料相談（30分）", "貴社の課題・ニーズをヒアリングします"),
    ("3", "ご契約（電子契約）", "クラウドサインによる電子契約で即手続き可"),
    ("4", "利用開始", "契約翌日からチャットで依頼開始できます"),
]):
    y = Inches(2.0 + i * 1.18)
    add_rect(slide17, Inches(0.9), y, Inches(6.0), Inches(1.0), TEAL_BG)
    add_rect(slide17, Inches(0.98), y + Inches(0.22), Inches(0.5), Inches(0.5), TEAL)
    add_text(slide17, num, Inches(0.98), y + Inches(0.2), Inches(0.5), Inches(0.5), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide17, h4, Inches(1.6), y + Inches(0.1), Inches(3.5), Inches(0.4), font_size=16, bold=True, color=DARK)
    add_text(slide17, p_text, Inches(1.6), y + Inches(0.52), Inches(5.1), Inches(0.38), font_size=13, color=GRAY)

add_rect(slide17, Inches(7.5), Inches(1.9), Inches(5.4), Inches(4.85), TEAL)
add_text(slide17, "★", Inches(7.5), Inches(2.25), Inches(5.4), Inches(0.7), font_size=32, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide17, "今なら\n[初回1ヶ月半額]\nキャンペーン実施中",
         Inches(7.8), Inches(2.95), Inches(4.8), Inches(1.9), font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide17, "まずはお気軽にご相談ください。\n御社の課題に合わせた最適なプランをご提案します。",
         Inches(7.8), Inches(4.95), Inches(4.8), Inches(1.2), font_size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ==================== SLIDE 18: FAQ ====================
slide18 = prs.slides.add_slide(blank_layout)
add_rect(slide18, 0, 0, W, H, WHITE)
slide_header(slide18, "Support", "サポート・FAQ")
slide_footer(slide18)

faqs = [
    ("修正回数に制限はありますか？", "A. 原則無制限です。軽微な変更は1営業日〜で対応します。"),
    ("途中解約はできますか？", "A. 更新月の[1]ヶ月前告知で解約可能です。違約金はありません。"),
    ("チャットツールは何に対応していますか？", "A. Chatwork, Slack, Teams等に対応しています。"),
    ("複数案件を同時に依頼できますか？", "A. 原則1案件ずつとなります。同時進行をご希望の場合は2口契約などをご検討ください。"),
    ("使い切れなかったチケットはどうなりますか？", "A. 翌月への繰り越しは[1ヶ月分]まで可能です。詳細はご契約時にご確認ください。"),
    ("支払い方法は？", "A. 月額払い・一括払いを選択可能です。一括払いは5%割引で請求書を発行します。"),
]
for i, (q, a) in enumerate(faqs):
    row = i // 2
    col = i % 2
    x = Inches(0.9 + col * 6.25)
    y = Inches(1.85 + row * 1.75)
    add_rect(slide18, x, y, Inches(6.0), Inches(1.6), TEAL_BG)
    add_rect(slide18, x + Inches(0.22), y + Inches(0.22), Inches(0.32), Inches(0.32), TEAL)
    add_text(slide18, "Q", x + Inches(0.22), y + Inches(0.2), Inches(0.32), Inches(0.32), font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide18, q, x + Inches(0.65), y + Inches(0.15), Inches(5.15), Inches(0.4), font_size=14, bold=True, color=DARK)
    add_text(slide18, a, x + Inches(0.27), y + Inches(0.68), Inches(5.5), Inches(0.75), font_size=13, color=DARK)


# ==================== SLIDE 19: WHY NOW ====================
slide19 = prs.slides.add_slide(blank_layout)
add_rect(slide19, 0, 0, W, H, TEAL)
slide_footer(slide19, teal_color=False)

add_text(slide19, "Why Now? なぜ今か", Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.2),
         font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide19,
         "迷っている間にも、競合は動いています。\n\nクリエイティブの改善は、後回しにするほど機会損失が大きくなります。\n今すぐ体制を整え、事業成長のスピードを加速させましょう。",
         Inches(1.5), Inches(2.6), Inches(10.3), Inches(2.4), font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide19, Inches(3.3), Inches(5.35), Inches(6.7), Inches(0.9), RGBColor(0x4D, 0xE5, 0xD0))
add_text(slide19, "今すぐ無料相談を予約する →",
         Inches(3.3), Inches(5.4), Inches(6.7), Inches(0.8), font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ==================== SLIDE 20: CONTACT ====================
slide20 = prs.slides.add_slide(blank_layout)
add_rect(slide20, 0, 0, W, H, WHITE)
slide_header(slide20, "Contact", "お問い合わせ", "まずは無料相談（30分）で、貴社の課題をお聞かせください。")
slide_footer(slide20)

add_text(slide20, "デザインの力で、\nあなたのビジネスを\n加速させましょう。",
         Inches(0.9), Inches(2.1), Inches(6.2), Inches(2.0), font_size=34, bold=True, color=DARK)
add_text(slide20, "まずは30分の無料相談からお気軽にどうぞ。\n貴社の課題に合わせた最適なプランをご提案します。",
         Inches(0.9), Inches(4.3), Inches(6.2), Inches(1.0), font_size=17, color=GRAY)

add_rect(slide20, Inches(7.4), Inches(1.85), Inches(5.6), Inches(5.0), TEAL_BG)
add_text(slide20, "お問い合わせ先", Inches(7.4), Inches(2.15), Inches(5.6), Inches(0.5),
         font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

for i, (icon, text) in enumerate([
    ("📋", "[フォームURL]"),
    ("✉️", "info@sakazuki.co.jp"),
    ("🏢", "株式会社SAKAZUKI"),
]):
    y = Inches(2.85 + i * 0.95)
    add_rect(slide20, Inches(7.75), y, Inches(0.46), Inches(0.46), TEAL)
    add_text(slide20, icon, Inches(7.75), y, Inches(0.46), Inches(0.46), font_size=16, align=PP_ALIGN.CENTER)
    add_text(slide20, text, Inches(8.32), y + Inches(0.05), Inches(4.4), Inches(0.4), font_size=16, color=DARK)

add_rect(slide20, Inches(7.7), Inches(5.65), Inches(4.9), Inches(0.03), RGBColor(0xE0, 0xE0, 0xE0))
add_text(slide20, "株式会社SAKAZUKI", Inches(7.4), Inches(5.85), Inches(5.6), Inches(0.45),
         font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)


# Save
output_path = "/Users/yoshidayuya/Desktop/claude/git/pass-presentation/PASS_presentation.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
